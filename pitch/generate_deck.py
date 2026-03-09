#!/usr/bin/env python3
"""
Agentic Agency Pitch Deck Generator

Generates a PPTX presentation matching the brutalist visual identity
from the Agentic Agency website.

Usage:
    python generate_deck.py
    python generate_deck.py --output custom_name.pptx
"""

import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import yaml


# =============================================================================
# VISUAL IDENTITY CONSTANTS
# =============================================================================

# Core Palette
COLOR_CEMENT = RGBColor(0xE6, 0xE6, 0xE1)  # #E6E6E1 - slide backgrounds
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)   # #000000 - headlines, text
COLOR_DARK = RGBColor(0x11, 0x11, 0x11)    # #111111 - accent slides, code
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF - contrast elements

# Typography (Space Grotesk preferred, Arial as fallback)
FONT_HEADLINE = "Space Grotesk"
FONT_BODY = "Space Grotesk"
FONT_MONO = "SF Mono"
FONT_FALLBACK = "Arial"

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Margins
MARGIN_LEFT = Inches(0.75)
MARGIN_RIGHT = Inches(0.75)
MARGIN_TOP = Inches(0.6)
MARGIN_BOTTOM = Inches(0.5)

# Content area
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_slide_background(slide, color):
    """Set solid color background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_name=FONT_HEADLINE, font_size=Pt(24),
                 bold=False, color=COLOR_BLACK, align=PP_ALIGN.LEFT,
                 vertical_anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.auto_size = None

    p = frame.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

    # Set vertical anchor
    textbox.text_frame.anchor = vertical_anchor

    return textbox


def add_multiline_text(slide, left, top, width, height, lines,
                       font_name=FONT_BODY, font_size=Pt(18),
                       color=COLOR_BLACK, line_spacing=1.5,
                       bullet=False):
    """Add multiple lines/paragraphs of text."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()

        if bullet:
            p.text = f"• {line}"
        else:
            p.text = line

        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = color
        p.line_spacing = line_spacing

    return textbox


def add_rectangle(slide, left, top, width, height, fill_color,
                  border_color=None, border_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    return shape


def add_badge(slide, left, top, text, bg_color=COLOR_BLACK,
              text_color=COLOR_WHITE):
    """Add a badge/tag element."""
    # Calculate width based on text length
    width = Inches(len(text) * 0.12 + 0.4)
    height = Inches(0.35)

    # Add background rectangle
    badge = add_rectangle(slide, left, top, width, height, bg_color)

    # Add text on top
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = False
    frame.anchor = MSO_ANCHOR.MIDDLE

    p = frame.paragraphs[0]
    p.text = text
    p.font.name = FONT_HEADLINE
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    return badge, textbox


# =============================================================================
# SLIDE TEMPLATES
# =============================================================================

def create_title_slide(prs, content):
    """
    Title Slide: Full-bleed cement background, large uppercase headline
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CEMENT)

    # Main headline - large, centered, uppercase
    headline = content.get('headline', 'THE AGENTIC AGENCY')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=Inches(2.2),
        width=CONTENT_WIDTH,
        height=Inches(1.5),
        text=headline.upper(),
        font_size=Pt(72),
        bold=True,
        color=COLOR_BLACK,
        align=PP_ALIGN.CENTER
    )

    # Subhead
    subhead = content.get('subhead', '')
    if subhead:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(3.8),
            width=CONTENT_WIDTH,
            height=Inches(0.6),
            text=subhead.upper(),
            font_size=Pt(28),
            bold=True,
            color=COLOR_BLACK,
            align=PP_ALIGN.CENTER
        )

    # Footer (date/venue)
    footer = content.get('footer', '')
    if footer:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(6.5),
            width=CONTENT_WIDTH,
            height=Inches(0.4),
            text=footer,
            font_size=Pt(14),
            color=COLOR_BLACK,
            align=PP_ALIGN.CENTER
        )

    return slide


def create_statement_slide(prs, content):
    """
    Statement Slide: Single powerful statement, centered
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CEMENT)

    # Optional headline at top
    headline = content.get('headline', '')
    if headline:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(0.8),
            width=CONTENT_WIDTH,
            height=Inches(0.6),
            text=headline.upper(),
            font_size=Pt(14),
            bold=True,
            color=COLOR_BLACK,
            align=PP_ALIGN.LEFT
        )

    # Main statement - centered, impactful
    statement = content.get('statement', '')
    add_text_box(
        slide,
        left=Inches(1.2),
        top=Inches(2.0),
        width=Inches(11),
        height=Inches(2.5),
        text=statement.strip(),
        font_size=Pt(32),
        bold=False,
        color=COLOR_BLACK,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE
    )

    # Emphasis line
    emphasis = content.get('emphasis', '')
    if emphasis:
        add_text_box(
            slide,
            left=Inches(1.2),
            top=Inches(4.8),
            width=Inches(11),
            height=Inches(0.8),
            text=emphasis,
            font_size=Pt(36),
            bold=True,
            color=COLOR_BLACK,
            align=PP_ALIGN.CENTER
        )

    return slide


def add_numbered_block(slide, number, text, left, top, width,
                        num_color=COLOR_BLACK, text_color=COLOR_BLACK,
                        show_line=True):
    """
    Add a numbered content block with large typographic number.
    Brutalist style: big number, thick underline, text beside it.
    """
    num_width = Inches(1.0)
    text_left = left + num_width + Inches(0.3)
    text_width = width - num_width - Inches(0.3)
    block_height = Inches(1.1)

    # Large number
    add_text_box(
        slide,
        left=left,
        top=top,
        width=num_width,
        height=Inches(0.9),
        text=f"{number:02d}",
        font_size=Pt(48),
        bold=True,
        color=num_color,
        align=PP_ALIGN.LEFT,
        vertical_anchor=MSO_ANCHOR.TOP
    )

    # Thick underline beneath number
    if show_line:
        add_rectangle(
            slide,
            left=left,
            top=top + Inches(0.75),
            width=Inches(0.7),
            height=Pt(4),
            fill_color=num_color
        )

    # Content text
    add_text_box(
        slide,
        left=text_left,
        top=top + Inches(0.05),
        width=text_width,
        height=block_height,
        text=text,
        font_size=Pt(22),
        bold=False,
        color=text_color,
        align=PP_ALIGN.LEFT,
        vertical_anchor=MSO_ANCHOR.TOP
    )

    return block_height


def add_content_card(slide, text, left, top, width, height,
                     bg_color=COLOR_WHITE, border_color=COLOR_BLACK,
                     text_color=COLOR_BLACK, border_width=Pt(3)):
    """
    Add a bordered card with content text.
    Brutalist style: thick black border, clean typography.
    """
    # Card background with thick border
    add_rectangle(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        fill_color=bg_color,
        border_color=border_color,
        border_width=border_width
    )

    # Content text with padding
    padding = Inches(0.2)
    add_text_box(
        slide,
        left=left + padding,
        top=top + padding,
        width=width - (padding * 2),
        height=height - (padding * 2),
        text=text,
        font_size=Pt(18),
        bold=False,
        color=text_color,
        align=PP_ALIGN.LEFT,
        vertical_anchor=MSO_ANCHOR.MIDDLE
    )


def create_content_slide(prs, content):
    """
    Content Slide: Brutalist numbered blocks or card grid layout.
    Uses large typographic numbers as design elements.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CEMENT)

    # Headline with thick underline
    headline = content.get('headline', '')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=MARGIN_TOP,
        width=CONTENT_WIDTH,
        height=Inches(0.8),
        text=headline.upper(),
        font_size=Pt(36),
        bold=True,
        color=COLOR_BLACK,
        align=PP_ALIGN.LEFT
    )

    # Thick line under headline
    add_rectangle(
        slide,
        left=MARGIN_LEFT,
        top=Inches(1.25),
        width=Inches(2.5),
        height=Pt(4),
        fill_color=COLOR_BLACK
    )

    # Content items
    bullets = content.get('bullets', [])
    layout = content.get('layout', 'numbered')  # 'numbered', 'cards', 'grid'

    if bullets:
        if layout == 'cards' or len(bullets) == 3:
            # Card grid layout - good for 3-4 items
            num_items = len(bullets)
            if num_items <= 2:
                # Side by side
                card_width = Inches(5.5)
                card_height = Inches(3.5)
                gap = Inches(0.5)
                start_left = MARGIN_LEFT
                for i, text in enumerate(bullets):
                    add_content_card(
                        slide,
                        text=text,
                        left=start_left + (i * (card_width + gap)),
                        top=Inches(1.8),
                        width=card_width,
                        height=card_height
                    )
            elif num_items == 3:
                # 3-column layout
                card_width = Inches(3.7)
                card_height = Inches(4.2)
                gap = Inches(0.35)
                start_left = MARGIN_LEFT
                for i, text in enumerate(bullets):
                    # Add number badge on card
                    card_left = start_left + (i * (card_width + gap))
                    card_top = Inches(1.8)

                    add_content_card(
                        slide,
                        text="",  # Empty, we'll add custom content
                        left=card_left,
                        top=card_top,
                        width=card_width,
                        height=card_height
                    )

                    # Large number in card
                    add_text_box(
                        slide,
                        left=card_left + Inches(0.25),
                        top=card_top + Inches(0.2),
                        width=Inches(0.8),
                        height=Inches(0.7),
                        text=f"{i+1:02d}",
                        font_size=Pt(32),
                        bold=True,
                        color=COLOR_BLACK,
                        align=PP_ALIGN.LEFT
                    )

                    # Thick underline
                    add_rectangle(
                        slide,
                        left=card_left + Inches(0.25),
                        top=card_top + Inches(0.85),
                        width=Inches(0.5),
                        height=Pt(3),
                        fill_color=COLOR_BLACK
                    )

                    # Text content
                    add_text_box(
                        slide,
                        left=card_left + Inches(0.25),
                        top=card_top + Inches(1.1),
                        width=card_width - Inches(0.5),
                        height=Inches(2.8),
                        text=text,
                        font_size=Pt(18),
                        bold=False,
                        color=COLOR_BLACK,
                        align=PP_ALIGN.LEFT
                    )
            else:
                # 2x2 grid for 4 items
                card_width = Inches(5.5)
                card_height = Inches(2.2)
                gap_x = Inches(0.5)
                gap_y = Inches(0.4)
                for i, text in enumerate(bullets[:4]):
                    row = i // 2
                    col = i % 2
                    card_left = MARGIN_LEFT + (col * (card_width + gap_x))
                    card_top = Inches(1.8) + (row * (card_height + gap_y))

                    add_content_card(
                        slide,
                        text="",
                        left=card_left,
                        top=card_top,
                        width=card_width,
                        height=card_height
                    )

                    # Number
                    add_text_box(
                        slide,
                        left=card_left + Inches(0.2),
                        top=card_top + Inches(0.15),
                        width=Inches(0.6),
                        height=Inches(0.5),
                        text=f"{i+1:02d}",
                        font_size=Pt(24),
                        bold=True,
                        color=COLOR_BLACK,
                        align=PP_ALIGN.LEFT
                    )

                    # Text
                    add_text_box(
                        slide,
                        left=card_left + Inches(0.9),
                        top=card_top + Inches(0.2),
                        width=card_width - Inches(1.1),
                        height=card_height - Inches(0.4),
                        text=text,
                        font_size=Pt(18),
                        bold=False,
                        color=COLOR_BLACK,
                        align=PP_ALIGN.LEFT,
                        vertical_anchor=MSO_ANCHOR.MIDDLE
                    )

        else:
            # Numbered blocks layout (default) - stacked with large numbers
            start_top = Inches(1.7)
            block_spacing = Inches(1.25)

            for i, text in enumerate(bullets):
                add_numbered_block(
                    slide,
                    number=i + 1,
                    text=text,
                    left=MARGIN_LEFT,
                    top=start_top + (i * block_spacing),
                    width=CONTENT_WIDTH
                )

    # Optional note at bottom
    note = content.get('note', '')
    if note:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(6.4),
            width=CONTENT_WIDTH,
            height=Inches(0.5),
            text=note,
            font_size=Pt(16),
            bold=True,
            color=COLOR_BLACK,
            align=PP_ALIGN.LEFT
        )

    return slide


def create_demo_anchor_slide(prs, content):
    """
    Demo Anchor Slide: Dark background, badge, key takeaway
    Used as visual anchors before/after live demos
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_DARK)

    # Badge (e.g., "LIVE DEMO" or "DEMO INSIGHT")
    badge_text = content.get('badge', 'DEMO')
    add_badge(
        slide,
        left=MARGIN_LEFT,
        top=Inches(1.5),
        text=badge_text,
        bg_color=COLOR_WHITE,
        text_color=COLOR_BLACK
    )

    # Main headline
    headline = content.get('headline', '')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=Inches(2.2),
        width=CONTENT_WIDTH,
        height=Inches(1.2),
        text=headline.upper(),
        font_size=Pt(52),
        bold=True,
        color=COLOR_WHITE,
        align=PP_ALIGN.LEFT
    )

    # Subhead
    subhead = content.get('subhead', '')
    if subhead:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(3.6),
            width=CONTENT_WIDTH,
            height=Inches(1.0),
            text=subhead,
            font_size=Pt(24),
            color=COLOR_WHITE,
            align=PP_ALIGN.LEFT
        )

    # Note (smaller, at bottom)
    note = content.get('note', '')
    if note:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(5.5),
            width=CONTENT_WIDTH,
            height=Inches(0.5),
            text=note,
            font_size=Pt(16),
            color=RGBColor(0x99, 0x99, 0x99),  # Muted gray
            align=PP_ALIGN.LEFT
        )

    # Transition indicator
    transition = content.get('transition', '')
    if transition:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(6.5),
            width=CONTENT_WIDTH,
            height=Inches(0.4),
            text=transition,
            font_size=Pt(14),
            bold=True,
            color=COLOR_WHITE,
            align=PP_ALIGN.LEFT
        )

    return slide


def create_value_promise_slide(prs, content):
    """
    Value Promise Slide: 2x3 grid of outcomes.
    Dark background with numbered cards showing what attendees will learn.
    Based on website's "What You'll Walk Away With" section.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_DARK)

    # Headline
    headline = content.get('headline', "WHAT YOU'LL WALK AWAY WITH")
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=Inches(0.4),
        width=CONTENT_WIDTH,
        height=Inches(0.7),
        text=headline.upper(),
        font_size=Pt(32),
        bold=True,
        color=COLOR_WHITE,
        align=PP_ALIGN.LEFT
    )

    # Subhead
    subhead = content.get('subhead', '')
    if subhead:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(1.0),
            width=CONTENT_WIDTH,
            height=Inches(0.4),
            text=subhead,
            font_size=Pt(16),
            color=RGBColor(0x99, 0x99, 0x99),
            align=PP_ALIGN.LEFT
        )

    # Outcomes grid (2x3)
    outcomes = content.get('outcomes', [])
    if outcomes:
        card_width = Inches(5.8)
        card_height = Inches(1.6)
        gap_x = Inches(0.4)
        gap_y = Inches(0.25)
        start_top = Inches(1.55)

        for i, outcome in enumerate(outcomes[:6]):
            row = i // 2
            col = i % 2
            card_left = MARGIN_LEFT + (col * (card_width + gap_x))
            card_top = start_top + (row * (card_height + gap_y))

            # Card background with subtle border
            add_rectangle(
                slide,
                left=card_left,
                top=card_top,
                width=card_width,
                height=card_height,
                fill_color=RGBColor(0x1A, 0x1A, 0x1A),
                border_color=RGBColor(0x33, 0x33, 0x33),
                border_width=Pt(2)
            )

            # Number circle
            num_size = Inches(0.35)
            add_rectangle(
                slide,
                left=card_left + Inches(0.2),
                top=card_top + Inches(0.2),
                width=num_size,
                height=num_size,
                fill_color=COLOR_CEMENT,
                border_color=None
            )
            add_text_box(
                slide,
                left=card_left + Inches(0.2),
                top=card_top + Inches(0.2),
                width=num_size,
                height=num_size,
                text=str(i + 1),
                font_size=Pt(14),
                bold=True,
                color=COLOR_BLACK,
                align=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE
            )

            # Title (bold part)
            title = outcome.get('title', '')
            add_text_box(
                slide,
                left=card_left + Inches(0.65),
                top=card_top + Inches(0.15),
                width=card_width - Inches(0.85),
                height=Inches(0.7),
                text=title,
                font_size=Pt(14),
                bold=True,
                color=COLOR_WHITE,
                align=PP_ALIGN.LEFT
            )

            # Benefit (lighter text)
            benefit = outcome.get('benefit', '')
            add_text_box(
                slide,
                left=card_left + Inches(0.65),
                top=card_top + Inches(0.85),
                width=card_width - Inches(0.85),
                height=Inches(0.6),
                text=benefit,
                font_size=Pt(12),
                color=RGBColor(0x99, 0x99, 0x99),
                align=PP_ALIGN.LEFT
            )

    return slide


def create_offer_slide(prs, content):
    """
    Offer Slide: Workshop/product offering with pricing.
    Brutalist style: bold typography, thick borders, clear hierarchy.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CEMENT)

    # Headline
    headline = content.get('headline', '')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=MARGIN_TOP,
        width=CONTENT_WIDTH,
        height=Inches(0.9),
        text=headline.upper(),
        font_size=Pt(48),
        bold=True,
        color=COLOR_BLACK,
        align=PP_ALIGN.LEFT
    )

    # Subhead badge
    subhead = content.get('subhead', '')
    if subhead:
        badge_width = Inches(len(subhead) * 0.13 + 0.5)
        add_rectangle(
            slide,
            left=MARGIN_LEFT,
            top=Inches(1.35),
            width=badge_width,
            height=Inches(0.4),
            fill_color=COLOR_BLACK
        )
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(1.35),
            width=badge_width,
            height=Inches(0.4),
            text=subhead.upper(),
            font_size=Pt(14),
            bold=True,
            color=COLOR_WHITE,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE
        )

    # Features as horizontal bars with left accent
    bullets = content.get('bullets', [])
    if bullets:
        bar_top = Inches(2.1)
        bar_height = Inches(0.7)
        bar_spacing = Inches(0.85)
        accent_width = Pt(6)

        for i, text in enumerate(bullets):
            current_top = bar_top + (i * bar_spacing)

            # Thick left accent bar
            add_rectangle(
                slide,
                left=MARGIN_LEFT,
                top=current_top,
                width=accent_width,
                height=bar_height,
                fill_color=COLOR_BLACK
            )

            # Text
            add_text_box(
                slide,
                left=MARGIN_LEFT + Inches(0.25),
                top=current_top,
                width=Inches(6.5),
                height=bar_height,
                text=text,
                font_size=Pt(20),
                bold=False,
                color=COLOR_BLACK,
                align=PP_ALIGN.LEFT,
                vertical_anchor=MSO_ANCHOR.MIDDLE
            )

    # Pricing cards - larger, more prominent
    pricing = content.get('pricing', [])
    if pricing:
        card_width = Inches(4.2)
        card_height = Inches(2.2)
        card_left = Inches(8.2)

        for i, price_item in enumerate(pricing):
            card_top = Inches(1.9 + i * 2.45)

            # Card background with thick border
            add_rectangle(
                slide,
                left=card_left,
                top=card_top,
                width=card_width,
                height=card_height,
                fill_color=COLOR_WHITE,
                border_color=COLOR_BLACK,
                border_width=Pt(4)
            )

            # Label
            add_text_box(
                slide,
                left=card_left + Inches(0.25),
                top=card_top + Inches(0.15),
                width=card_width - Inches(0.5),
                height=Inches(0.35),
                text=price_item.get('label', '').upper(),
                font_size=Pt(12),
                bold=True,
                color=COLOR_BLACK,
                align=PP_ALIGN.LEFT
            )

            # Price - large and bold
            add_text_box(
                slide,
                left=card_left + Inches(0.25),
                top=card_top + Inches(0.5),
                width=card_width - Inches(0.5),
                height=Inches(0.55),
                text=price_item.get('price', ''),
                font_size=Pt(28),
                bold=True,
                color=COLOR_BLACK,
                align=PP_ALIGN.LEFT
            )

            # Detail/terms - smaller text
            detail = price_item.get('detail', '')
            if detail:
                add_text_box(
                    slide,
                    left=card_left + Inches(0.25),
                    top=card_top + Inches(1.15),
                    width=card_width - Inches(0.5),
                    height=Inches(0.8),
                    text=detail,
                    font_size=Pt(11),
                    color=RGBColor(0x55, 0x55, 0x55),
                    align=PP_ALIGN.LEFT
                )

    return slide


def create_cta_slide(prs, content):
    """
    CTA Slide: Clear call-to-action with contact info
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_DARK)

    # Headline
    headline = content.get('headline', 'NEXT STEP')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=Inches(1.5),
        width=CONTENT_WIDTH,
        height=Inches(1.0),
        text=headline.upper(),
        font_size=Pt(56),
        bold=True,
        color=COLOR_WHITE,
        align=PP_ALIGN.CENTER
    )

    # CTA text
    cta = content.get('cta', '')
    if cta:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(2.8),
            width=CONTENT_WIDTH,
            height=Inches(0.8),
            text=cta,
            font_size=Pt(32),
            color=COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

    # Contact info
    contact = content.get('contact', [])
    if contact:
        add_multiline_text(
            slide,
            left=MARGIN_LEFT,
            top=Inches(4.2),
            width=CONTENT_WIDTH,
            height=Inches(1.5),
            lines=contact,
            font_size=Pt(24),
            color=COLOR_WHITE,
            line_spacing=1.5,
            bullet=False
        )
        # Center the contact textbox
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for p in shape.text_frame.paragraphs:
                    if any(c in p.text for c in contact):
                        p.alignment = PP_ALIGN.CENTER

    # Footer
    footer = content.get('footer', '')
    if footer:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(6.2),
            width=CONTENT_WIDTH,
            height=Inches(0.5),
            text=footer,
            font_size=Pt(20),
            bold=True,
            color=COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

    return slide


# =============================================================================
# DATA-DRIVEN SLIDE TYPES
# =============================================================================

def create_data_table_slide(prs, content):
    """
    Data Table Slide: Shows tabular data with clear headers.
    Brutalist style: thick borders, monospace numbers.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background based on theme
    theme = content.get('theme', 'light')
    if theme == 'dark':
        set_slide_background(slide, COLOR_DARK)
        text_color = COLOR_WHITE
        header_bg = COLOR_WHITE
        header_text = COLOR_BLACK
        cell_bg = RGBColor(0x1A, 0x1A, 0x1A)
        border_color = RGBColor(0x33, 0x33, 0x33)
    else:
        set_slide_background(slide, COLOR_CEMENT)
        text_color = COLOR_BLACK
        header_bg = COLOR_BLACK
        header_text = COLOR_WHITE
        cell_bg = COLOR_WHITE
        border_color = COLOR_BLACK

    # Headline
    headline = content.get('headline', '')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=Inches(0.4),
        width=CONTENT_WIDTH,
        height=Inches(0.6),
        text=headline.upper(),
        font_size=Pt(28),
        bold=True,
        color=text_color,
        align=PP_ALIGN.LEFT
    )

    # Subhead / assumption text
    subhead = content.get('subhead', '')
    if subhead:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(0.95),
            width=CONTENT_WIDTH,
            height=Inches(0.4),
            text=subhead,
            font_size=Pt(14),
            color=RGBColor(0x66, 0x66, 0x66) if theme == 'light' else RGBColor(0x99, 0x99, 0x99),
            align=PP_ALIGN.LEFT
        )

    # Table data
    table_data = content.get('table', {})
    headers = table_data.get('headers', [])
    rows = table_data.get('rows', [])

    if headers and rows:
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header

        # Calculate table dimensions
        table_top = Inches(1.5)
        table_height = Inches(4.8)
        row_height = table_height / num_rows

        # Create table
        table = slide.shapes.add_table(
            num_rows, num_cols,
            MARGIN_LEFT, table_top,
            CONTENT_WIDTH, table_height
        ).table

        # Style header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg

            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT_HEADLINE
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = header_text
            para.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

        # Style data rows
        for row_idx, row in enumerate(rows):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value) if value is not None else ''
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_bg

                para = cell.text_frame.paragraphs[0]
                para.font.name = FONT_MONO if col_idx > 0 else FONT_BODY
                para.font.size = Pt(10)
                para.font.color.rgb = text_color
                para.alignment = PP_ALIGN.RIGHT if col_idx > 0 else PP_ALIGN.LEFT

    # Assumptions box at bottom
    assumptions = content.get('assumptions', [])
    if assumptions:
        assumptions_text = "ASSUMPTIONS: " + " · ".join(assumptions)
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(6.5),
            width=CONTENT_WIDTH,
            height=Inches(0.5),
            text=assumptions_text,
            font_size=Pt(10),
            color=RGBColor(0x66, 0x66, 0x66) if theme == 'light' else RGBColor(0x88, 0x88, 0x88),
            align=PP_ALIGN.LEFT
        )

    return slide


def create_chart_slide(prs, content):
    """
    Chart Slide: Simple bar/line visualization.
    Uses rectangles to create bar charts (python-pptx chart API is complex).
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CEMENT)

    # Headline
    headline = content.get('headline', '')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=Inches(0.4),
        width=CONTENT_WIDTH,
        height=Inches(0.6),
        text=headline.upper(),
        font_size=Pt(28),
        bold=True,
        color=COLOR_BLACK,
        align=PP_ALIGN.LEFT
    )

    # Subhead
    subhead = content.get('subhead', '')
    if subhead:
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(0.95),
            width=CONTENT_WIDTH,
            height=Inches(0.4),
            text=subhead,
            font_size=Pt(14),
            color=RGBColor(0x66, 0x66, 0x66),
            align=PP_ALIGN.LEFT
        )

    # Chart data
    chart_data = content.get('chart', {})
    labels = chart_data.get('labels', [])
    values = chart_data.get('values', [])
    chart_type = chart_data.get('type', 'bar')

    if labels and values:
        # Chart area
        chart_left = MARGIN_LEFT + Inches(0.5)
        chart_top = Inches(1.8)
        chart_width = CONTENT_WIDTH - Inches(1)
        chart_height = Inches(4.0)

        # Calculate max value for scaling
        max_val = max(values)
        min_val = min(values)
        value_range = max_val - min_val if max_val != min_val else max_val

        # Draw baseline (for values that go negative)
        if min_val < 0:
            zero_y = chart_top + chart_height * (max_val / value_range)
            add_rectangle(
                slide,
                left=chart_left,
                top=zero_y,
                width=chart_width,
                height=Pt(2),
                fill_color=COLOR_BLACK
            )

        # Draw bars
        num_bars = len(values)
        bar_width = chart_width / (num_bars * 1.5)
        bar_gap = bar_width * 0.5

        for i, (label, value) in enumerate(zip(labels, values)):
            bar_left = chart_left + (i * (bar_width + bar_gap))

            # Calculate bar height and position
            if value >= 0:
                bar_height = (value / max_val) * chart_height * 0.9
                bar_top = chart_top + chart_height - bar_height - Inches(0.3)
                bar_color = COLOR_BLACK
            else:
                bar_height = abs(value / min_val) * chart_height * 0.3
                bar_top = chart_top + chart_height * 0.7
                bar_color = RGBColor(0x88, 0x00, 0x00)

            # Draw bar
            add_rectangle(
                slide,
                left=bar_left,
                top=bar_top,
                width=bar_width,
                height=bar_height,
                fill_color=bar_color
            )

            # Label below bar
            add_text_box(
                slide,
                left=bar_left - Inches(0.1),
                top=chart_top + chart_height + Inches(0.05),
                width=bar_width + Inches(0.2),
                height=Inches(0.3),
                text=str(label),
                font_size=Pt(8),
                color=COLOR_BLACK,
                align=PP_ALIGN.CENTER
            )

            # Value on bar (European number format: dots as thousands separator)
            formatted_value = f"{value:,.0f}".replace(",", ".")
            add_text_box(
                slide,
                left=bar_left - Inches(0.2),
                top=bar_top - Inches(0.25),
                width=bar_width + Inches(0.4),
                height=Inches(0.25),
                text=formatted_value,
                font_size=Pt(8),
                bold=True,
                color=COLOR_BLACK,
                align=PP_ALIGN.CENTER
            )

    # Key insight box
    insight = content.get('insight', '')
    if insight:
        add_rectangle(
            slide,
            left=MARGIN_LEFT,
            top=Inches(6.3),
            width=CONTENT_WIDTH,
            height=Inches(0.6),
            fill_color=COLOR_BLACK
        )
        add_text_box(
            slide,
            left=MARGIN_LEFT + Inches(0.2),
            top=Inches(6.35),
            width=CONTENT_WIDTH - Inches(0.4),
            height=Inches(0.5),
            text=insight,
            font_size=Pt(12),
            bold=True,
            color=COLOR_WHITE,
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_ANCHOR.MIDDLE
        )

    return slide


def create_matrix_slide(prs, content):
    """
    Matrix Slide: Who/What/When responsibility matrix.
    Clean grid with thick borders and clear ownership.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_CEMENT)

    # Headline
    headline = content.get('headline', '')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=Inches(0.4),
        width=CONTENT_WIDTH,
        height=Inches(0.6),
        text=headline.upper(),
        font_size=Pt(28),
        bold=True,
        color=COLOR_BLACK,
        align=PP_ALIGN.LEFT
    )

    # Matrix data
    matrix = content.get('matrix', {})
    headers = matrix.get('headers', ['WHO', 'WHAT', 'WHEN', 'STATUS'])
    rows = matrix.get('rows', [])

    if rows:
        num_cols = len(headers)
        num_rows = len(rows) + 1

        table_top = Inches(1.2)
        table_height = min(Inches(5.5), Inches(0.6 * num_rows))

        table = slide.shapes.add_table(
            num_rows, num_cols,
            MARGIN_LEFT, table_top,
            CONTENT_WIDTH, table_height
        ).table

        # Column widths (first column wider for WHO)
        col_widths = [Inches(2.0), Inches(5.5), Inches(2.0), Inches(2.3)]
        for i, width in enumerate(col_widths[:num_cols]):
            table.columns[i].width = width

        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header).upper()
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BLACK

            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT_HEADLINE
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = COLOR_WHITE
            para.alignment = PP_ALIGN.LEFT

        # Data rows
        for row_idx, row in enumerate(rows):
            for col_idx, value in enumerate(row[:num_cols]):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value) if value else ''
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE

                para = cell.text_frame.paragraphs[0]
                para.font.name = FONT_BODY
                para.font.size = Pt(10)
                para.font.bold = (col_idx == 0)  # Bold the WHO column
                para.font.color.rgb = COLOR_BLACK
                para.alignment = PP_ALIGN.LEFT

    return slide


def create_pyramid_headline_slide(prs, content):
    """
    Pyramid Headline Slide: Big conclusion statement with supporting points.
    Used for pyramid principle structure - conclusion first.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    theme = content.get('theme', 'dark')
    if theme == 'dark':
        set_slide_background(slide, COLOR_DARK)
        text_color = COLOR_WHITE
        muted_color = RGBColor(0x99, 0x99, 0x99)
    else:
        set_slide_background(slide, COLOR_CEMENT)
        text_color = COLOR_BLACK
        muted_color = RGBColor(0x66, 0x66, 0x66)

    # Main conclusion - large and centered
    headline = content.get('headline', '')
    add_text_box(
        slide,
        left=MARGIN_LEFT,
        top=Inches(1.5),
        width=CONTENT_WIDTH,
        height=Inches(2.0),
        text=headline.upper(),
        font_size=Pt(48),
        bold=True,
        color=text_color,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE
    )

    # Supporting statement
    subhead = content.get('subhead', '')
    if subhead:
        add_text_box(
            slide,
            left=Inches(1.5),
            top=Inches(3.8),
            width=Inches(10.333),
            height=Inches(1.0),
            text=subhead,
            font_size=Pt(20),
            color=muted_color,
            align=PP_ALIGN.CENTER
        )

    # Evidence bullets at bottom
    evidence = content.get('evidence', [])
    if evidence:
        evidence_text = " · ".join(evidence)
        add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(5.5),
            width=CONTENT_WIDTH,
            height=Inches(0.8),
            text=evidence_text,
            font_size=Pt(14),
            color=muted_color,
            align=PP_ALIGN.CENTER
        )

    return slide


# =============================================================================
# SLIDE TYPE DISPATCHER
# =============================================================================

SLIDE_CREATORS = {
    'title': create_title_slide,
    'statement': create_statement_slide,
    'content': create_content_slide,
    'demo_anchor': create_demo_anchor_slide,
    'value_promise': create_value_promise_slide,
    'offer': create_offer_slide,
    'cta': create_cta_slide,
    'data_table': create_data_table_slide,
    'chart': create_chart_slide,
    'matrix': create_matrix_slide,
    'pyramid': create_pyramid_headline_slide,
}


def create_slide(prs, slide_content):
    """Create a slide based on its type."""
    slide_type = slide_content.get('type', 'content')
    creator = SLIDE_CREATORS.get(slide_type, create_content_slide)
    return creator(prs, slide_content)


# =============================================================================
# MAIN GENERATOR
# =============================================================================

def load_content(content_path):
    """Load slide content from YAML file."""
    with open(content_path, 'r') as f:
        return yaml.safe_load(f)


def generate_deck(content_path, output_path):
    """Generate the full presentation."""
    # Load content
    content = load_content(content_path)

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Generate each slide
    slides_content = content.get('slides', [])
    for i, slide_content in enumerate(slides_content):
        slide_id = slide_content.get('id', f'slide_{i}')
        print(f"  Creating slide {i+1}: {slide_id}")
        create_slide(prs, slide_content)

    # Save presentation
    prs.save(output_path)
    print(f"\n✓ Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate Agentic Agency pitch deck"
    )
    parser.add_argument(
        '--content',
        default='content/slides.yaml',
        help='Path to slides content YAML file'
    )
    parser.add_argument(
        '--output',
        default='output/agentic_agency_pitch.pptx',
        help='Output path for the PPTX file'
    )
    args = parser.parse_args()

    # Resolve paths relative to script location
    script_dir = Path(__file__).parent
    content_path = script_dir / args.content
    output_path = script_dir / args.output

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)

    print("Generating Agentic Agency Pitch Deck")
    print("=" * 40)
    print(f"Content: {content_path}")
    print(f"Output:  {output_path}")
    print()

    generate_deck(content_path, output_path)


if __name__ == '__main__':
    main()
